import logging
import math
from pathlib import Path
from typing import Any

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Emu, Pt


logger = logging.getLogger(__name__)

EMU_PER_INCH = 914400
MAX_EMU = 20116800  # 22 дюйма — жёсткий лимит python-pptx / PowerPoint


class PPTXGenerator:
    """
    Генератор PPTX из OCR-результатов.

    Важно:
    - изображение НЕ масштабируется под слайд;
    - слайд делается под физический размер изображения;
    - OCR-координаты переводятся из px в EMU через image_dpi;
    - для PDF image_dpi должен совпадать с dpi в PaddleOCRProcessor._pdf_to_images().
    """

    def __init__(
        self,
        image_dpi: int = 200,
        font_name: str = "Calibri",
        font_color: tuple[int, int, int] = (0, 0, 0),
        min_font_pt: int = 7,
        max_font_pt: int = 38,
        bbox_x_offset_px: float = 0.0,
        bbox_y_offset_px: float = 0.0,
        vertical_center_text: bool = True,
    ) -> None:
        if image_dpi <= 0:
            raise ValueError("image_dpi должен быть положительным")

        self.image_dpi = image_dpi

        self.font_name = font_name
        self.font_color = font_color
        self.min_font_pt = min_font_pt
        self.max_font_pt = max_font_pt

        self.bbox_x_offset_px = bbox_x_offset_px
        self.bbox_y_offset_px = bbox_y_offset_px
        self.vertical_center_text = vertical_center_text

        self.emu_per_px = EMU_PER_INCH / image_dpi

        self.slide_width_emu: int = 0
        self.slide_height_emu: int = 0

    @staticmethod
    def _clamp_emu(value: float) -> int:
        return int(max(0, min(value, MAX_EMU)))

    def _safe_emu(self, value: float, label: str = "") -> Emu:
        clamped = self._clamp_emu(value)

        if clamped != int(value) and label:
            logger.warning(
                "⚠️ Значение %s ограничено: %.0f → %d EMU",
                label,
                value,
                clamped,
            )

        return Emu(clamped)

    def _px_to_emu(self, value_px: float) -> float:
        """
        Переводит пиксели исходного изображения в EMU без fit-resize.
        """
        return value_px * self.emu_per_px

    def _parse_img_size(self, img_size: Any) -> tuple[int, int]:
        if not img_size:
            return 0, 0

        try:
            img_w, img_h = img_size
            return int(img_w), int(img_h)
        except Exception as exc:
            logger.warning("⚠️ Некорректный img_size=%r: %s", img_size, exc)
            return 0, 0

    def _parse_bbox(self, box: Any) -> tuple[int, int, int, int]:
        """
        Поддерживает:
        - [[x1, y1], [x2, y2], ...]
        - [x1, y1, x2, y2, ...]
        """
        try:
            if not isinstance(box, (list, tuple)) or len(box) == 0:
                return 0, 0, 0, 0

            points: list[tuple[float, float]] = []

            if isinstance(box[0], (list, tuple)):
                for point in box:
                    if len(point) >= 2:
                        points.append((float(point[0]), float(point[1])))
            else:
                usable_len = len(box) - (len(box) % 2)

                if usable_len < 4:
                    return 0, 0, 0, 0

                for i in range(0, usable_len, 2):
                    points.append((float(box[i]), float(box[i + 1])))

            if not points:
                return 0, 0, 0, 0

            xs = [p[0] for p in points]
            ys = [p[1] for p in points]

            left = max(0, int(min(xs)))
            top = max(0, int(min(ys)))
            right = max(xs)
            bottom = max(ys)

            width = max(0, int(right - left))
            height = max(0, int(bottom - top))

            return left, top, width, height

        except Exception as exc:
            logger.debug("⚠️ Не удалось распарсить bbox=%r: %s", box, exc)
            return 0, 0, 0, 0

    def _calculate_font_size(
        self,
        box_h_emu: int,
        text: str,
        box_w_emu: int,
    ) -> Pt:
        if box_h_emu <= 0 or box_w_emu <= 0:
            return Pt(12)

        text = text or ""

        if not text.strip():
            return Pt(12)

        h_pt = (box_h_emu / EMU_PER_INCH) * 72
        w_pt = (box_w_emu / EMU_PER_INCH) * 72

        avg_char_width_pt = 5.5
        chars_per_line = max(1, w_pt / avg_char_width_pt)
        lines_needed = max(1, math.ceil(len(text) / chars_per_line))

        line_height_pt = h_pt / lines_needed
        font_pt = int(line_height_pt * 0.80)

        font_pt = max(self.min_font_pt, min(self.max_font_pt, font_pt))

        return Pt(font_pt)

    def _clip_rect_to_slide(
        self,
        left: float,
        top: float,
        width: float,
        height: float,
    ) -> tuple[float, float, float, float]:
        right = left + width
        bottom = top + height

        clipped_left = max(0.0, left)
        clipped_top = max(0.0, top)
        clipped_right = min(float(self.slide_width_emu), right)
        clipped_bottom = min(float(self.slide_height_emu), bottom)

        clipped_width = clipped_right - clipped_left
        clipped_height = clipped_bottom - clipped_top

        return clipped_left, clipped_top, clipped_width, clipped_height

    def _find_first_valid_image_size(
        self,
        ocr_results: list[dict],
    ) -> tuple[int, int]:
        """
        PowerPoint не поддерживает разный размер слайда внутри одной презентации.
        Поэтому размер презентации берём по первой валидной странице.
        """
        for page in ocr_results:
            if not isinstance(page, dict):
                continue

            img_w, img_h = self._parse_img_size(page.get("img_size"))

            if img_w > 1 and img_h > 1:
                return img_w, img_h

        return 0, 0

    def _setup_slide_size_from_first_image(
        self,
        prs: Presentation,
        ocr_results: list[dict],
    ) -> None:
        img_w, img_h = self._find_first_valid_image_size(ocr_results)

        if img_w <= 1 or img_h <= 1:
            raise ValueError("Не удалось определить размер первой валидной страницы")

        raw_slide_width = self._px_to_emu(img_w)
        raw_slide_height = self._px_to_emu(img_h)

        self.slide_width_emu = self._clamp_emu(raw_slide_width)
        self.slide_height_emu = self._clamp_emu(raw_slide_height)

        if self.slide_width_emu != int(raw_slide_width):
            logger.warning(
                "⚠️ Ширина слайда превышает лимит PowerPoint: %.0f → %d EMU. "
                "Оригинальный размер полностью сохранить нельзя.",
                raw_slide_width,
                self.slide_width_emu,
            )

        if self.slide_height_emu != int(raw_slide_height):
            logger.warning(
                "⚠️ Высота слайда превышает лимит PowerPoint: %.0f → %d EMU. "
                "Оригинальный размер полностью сохранить нельзя.",
                raw_slide_height,
                self.slide_height_emu,
            )

        prs.slide_width = self.slide_width_emu
        prs.slide_height = self.slide_height_emu

        logger.info(
            "Размер слайда установлен по изображению: %d x %d px, dpi=%d, "
            "slide=%d x %d EMU",
            img_w,
            img_h,
            self.image_dpi,
            self.slide_width_emu,
            self.slide_height_emu,
        )

    def _add_ocr_textbox(
        self,
        slide: Any,
        left: float,
        top: float,
        width: float,
        height: float,
        text: str,
    ) -> None:
        if width <= 0 or height <= 0:
            return

        textbox = slide.shapes.add_textbox(
            self._safe_emu(left, "tx_left"),
            self._safe_emu(top, "tx_top"),
            self._safe_emu(width, "tx_width"),
            self._safe_emu(height, "tx_height"),
        )

        text_frame = textbox.text_frame
        text_frame.word_wrap = True
        text_frame.margin_left = Emu(0)
        text_frame.margin_right = Emu(0)
        text_frame.margin_top = Emu(0)
        text_frame.margin_bottom = Emu(0)

        if self.vertical_center_text:
            text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        else:
            text_frame.vertical_anchor = MSO_ANCHOR.TOP

        paragraph = text_frame.paragraphs[0]
        paragraph.text = text
        paragraph.alignment = PP_ALIGN.LEFT
        paragraph.space_before = Pt(0)
        paragraph.space_after = Pt(0)

        font_size = self._calculate_font_size(
            int(height),
            text,
            int(width),
        )

        paragraph.font.size = font_size
        paragraph.line_spacing = Pt(font_size.pt * 1.05)
        paragraph.font.name = self.font_name
        paragraph.font.color.rgb = RGBColor(*self.font_color)

    def generate(self, ocr_results: list[dict], output_path: str) -> None:
        """
        Генерирует PPTX без масштабирования изображения под слайд.

        Схема:
        - размер слайда = размер первой картинки в физических единицах;
        - картинка вставляется в (0, 0);
        - размер картинки = img_w / dpi x img_h / dpi;
        - OCR-боксы переводятся px → EMU тем же коэффициентом.
        """
        if not ocr_results:
            raise ValueError("Пустой список OCR-результатов")

        prs = Presentation()
        self._setup_slide_size_from_first_image(prs, ocr_results)

        blank_layout = prs.slide_layouts[6]

        for page_index, page in enumerate(ocr_results, start=1):
            if not isinstance(page, dict):
                logger.warning(
                    "⚠️ Страница #%d пропущена: ожидался dict, получено %s",
                    page_index,
                    type(page).__name__,
                )
                continue

            img_path_raw = page.get("img_path")
            img_path = Path(img_path_raw) if img_path_raw else None

            img_w, img_h = self._parse_img_size(page.get("img_size"))

            if not img_path:
                logger.warning("⚠️ Страница #%d пропущена: отсутствует img_path", page_index)
                continue

            if not img_path.exists():
                logger.warning(
                    "⚠️ Страница #%d пропущена: файл изображения не найден: %s",
                    page_index,
                    img_path,
                )
                continue

            if img_w <= 1 or img_h <= 1:
                logger.warning(
                    "⚠️ Страница #%d пропущена: некорректный img_size=%r",
                    page_index,
                    page.get("img_size"),
                )
                continue

            slide = prs.slides.add_slide(blank_layout)

            # Главное изменение:
            # больше нет _fit_image_to_slide().
            # Изображение вставляется с оригинальным физическим размером.
            img_left = 0.0
            img_top = 0.0
            img_width = self._px_to_emu(img_w)
            img_height = self._px_to_emu(img_h)

            # Если страница больше лимита PowerPoint, физически сохранить
            # оригинальный размер нельзя. Ограничиваем, чтобы python-pptx не упал.
            img_width = min(img_width, MAX_EMU)
            img_height = min(img_height, MAX_EMU)

            try:
                slide.shapes.add_picture(
                    str(img_path),
                    self._safe_emu(img_left, "img_left"),
                    self._safe_emu(img_top, "img_top"),
                    self._safe_emu(img_width, "img_width"),
                    self._safe_emu(img_height, "img_height"),
                )
            except Exception as exc:
                logger.warning(
                    "⚠️ Не удалось добавить изображение на странице #%d: %s",
                    page_index,
                    exc,
                )
                continue

            texts = page.get("texts") or []

            if not isinstance(texts, list):
                logger.warning(
                    "⚠️ На странице #%d поле texts имеет некорректный тип: %s",
                    page_index,
                    type(texts).__name__,
                )
                continue

            for item_index, item in enumerate(texts, start=1):
                if not isinstance(item, dict):
                    logger.debug(
                        "⚠️ OCR-блок #%d на странице #%d пропущен: ожидался dict",
                        item_index,
                        page_index,
                    )
                    continue

                text = str(item.get("text") or "")

                if not text.strip():
                    continue

                ox, oy, ow, oh = self._parse_bbox(item.get("box"))

                # Снизил порог, чтобы не терять мелкий текст у краёв.
                if ow < 3 or oh < 3:
                    logger.debug(
                        "⚠️ Маленький bbox пропущен: page=%d item=%d text=%r bbox=%r parsed=%r",
                        page_index,
                        item_index,
                        text[:50],
                        item.get("box"),
                        (ox, oy, ow, oh),
                    )
                    continue

                # Координаты считаются напрямую из пикселей.
                # Никакого fit-resize относительно слайда больше нет.
                tx_left = self._px_to_emu(ox + self.bbox_x_offset_px)
                tx_top = self._px_to_emu(oy + self.bbox_y_offset_px)
                tx_width = self._px_to_emu(ow)
                tx_height = self._px_to_emu(oh)

                tx_left, tx_top, tx_width, tx_height = self._clip_rect_to_slide(
                    tx_left,
                    tx_top,
                    tx_width,
                    tx_height,
                )

                if tx_width <= 0 or tx_height <= 0:
                    continue

                self._add_ocr_textbox(
                    slide=slide,
                    left=tx_left,
                    top=tx_top,
                    width=tx_width,
                    height=tx_height,
                    text=text,
                )

        output = Path(output_path)

        if output.parent and not output.parent.exists():
            output.parent.mkdir(parents=True, exist_ok=True)

        prs.save(str(output))
        logger.info("✅ PPTX сохранён: %s", output)